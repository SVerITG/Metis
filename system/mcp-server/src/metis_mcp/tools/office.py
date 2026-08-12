"""office.py — read PowerPoint and Excel back INTO Metis (Keystone P5.2).

Metis could already write a .pptx. It could not read one, so anything the
researcher changed in PowerPoint was lost to Metis the moment they saved it — a
one-way integration, which is the same thing as no integration.

A DELIBERATE ASYMMETRY BETWEEN THE TWO FORMATS
    A .pptx is AUTHORED CONTENT: titles, bullets, speaker notes. Reading it whole
    is exactly what makes it useful, and it is the researcher's own prose.

    An .xlsx is DATA, and in this researcher's work that routinely means patient
    records, line-lists and screening exports. So the spreadsheet reader returns
    STRUCTURE ONLY — sheet names, dimensions, column headers, inferred types — and
    never cell values. Dumping the cells would drive a coach and horses through the
    "send code, not data" doctrine that the rest of Metis is built on, and it would
    do it through a file-reading convenience nobody would think to audit.

    If the researcher wants the values analysed, that is what the safe-analysis
    procedure is for: Metis writes a script, they run it locally, aggregates come
    back. This module makes that the only path, rather than the polite one.
"""
from __future__ import annotations

import datetime
import logging
from pathlib import Path

from mcp.types import TextContent

from metis_mcp.app_instance import app
from metis_mcp.config import paths
from metis_mcp.db import connect

log = logging.getLogger("metis.office")

MAX_SLIDE_CHARS = 4000
MAX_HEADER_COLS = 60


def _resolve(path: str) -> Path | None:
    p = Path(path).expanduser()
    if not p.is_absolute():
        p = paths.root / path
    return p if p.is_file() else None


def _read_pptx(p: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(p))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title, body = "", []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            if not title and shape == slide.shapes.title:
                title = text
            else:
                body.append(text)
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            pass
        slides.append({"number": i, "title": title or f"Slide {i}",
                       "body": "\n".join(body)[:MAX_SLIDE_CHARS], "notes": notes[:1000]})
    return {"kind": "pptx", "slides": slides,
            "slide_count": len(slides),
            "size": f"{prs.slide_width.inches:.2f}x{prs.slide_height.inches:.2f} in"}


def _read_xlsx(p: Path) -> dict:
    """Structure only — never cell values. See the module docstring."""
    import openpyxl

    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        headers, types = [], []
        try:
            rows = ws.iter_rows(min_row=1, max_row=2, max_col=MAX_HEADER_COLS,
                                values_only=True)
            first = next(rows, ()) or ()
            second = next(rows, ()) or ()
            headers = [str(h)[:60] for h in first if h is not None]
            # A single sample row is used ONLY to infer a type name, never reported.
            for v in second[:len(headers)]:
                types.append(type(v).__name__ if v is not None else "empty")
        except Exception:
            pass
        sheets.append({"name": ws.title,
                       "rows": ws.max_row, "columns": ws.max_column,
                       "headers": headers, "column_types": types})
    wb.close()
    return {"kind": "xlsx", "sheets": sheets, "sheet_count": len(sheets)}


@app.tool()
async def read_office_document(path: str) -> list[TextContent]:
    """Read a PowerPoint or Excel file into Metis.

    PowerPoint: returns the slides — titles, body text and speaker notes — because
    a deck is the researcher's own authored content.

    Excel: returns STRUCTURE ONLY — sheet names, row/column counts, column headers
    and inferred types. Never cell values. If you need the data analysed, use the
    safe-analysis procedure: Metis writes a script, the researcher runs it locally,
    and only aggregates come back.

    Args:
        path: Path to a .pptx or .xlsx file, absolute or relative to the Metis folder.
    """
    p = _resolve(path)
    if not p:
        return [TextContent(type="text", text=f"No file at '{path}'.")]
    suffix = p.suffix.lower()
    try:
        if suffix == ".pptx":
            data = _read_pptx(p)
            lines = [f"**{p.name}** — {data['slide_count']} slide(s), {data['size']}", ""]
            for s in data["slides"]:
                lines.append(f"### {s['number']}. {s['title']}")
                if s["body"]:
                    lines.append(s["body"])
                if s["notes"]:
                    lines.append(f"_Notes: {s['notes']}_")
                lines.append("")
            return [TextContent(type="text", text="\n".join(lines))]

        if suffix in (".xlsx", ".xlsm"):
            data = _read_xlsx(p)
            lines = [f"**{p.name}** — {data['sheet_count']} sheet(s). "
                     f"Structure only; no cell values are read.", ""]
            for sh in data["sheets"]:
                lines.append(f"### {sh['name']} — {sh['rows']} rows × {sh['columns']} columns")
                if sh["headers"]:
                    pairs = ", ".join(
                        f"{h} ({t})" for h, t in zip(sh["headers"], sh["column_types"] or [""] * len(sh["headers"]))
                    )
                    lines.append(f"Columns: {pairs}")
                lines.append("")
            lines.append("_To analyse the contents, use the safe-analysis procedure — "
                         "Metis writes a script you run locally, and only aggregates come back._")
            return [TextContent(type="text", text="\n".join(lines))]

        return [TextContent(type="text", text=f"'{p.suffix}' is not supported — .pptx or .xlsx only.")]
    except ImportError as exc:
        return [TextContent(type="text", text=f"The reader for {suffix} is not installed: {exc}")]
    except Exception as exc:
        return [TextContent(type="text", text=f"Could not read {p.name}: {type(exc).__name__}: {exc}")]


@app.tool()
async def ingest_office_document(path: str, project_id: str = "") -> list[TextContent]:
    """Read a PowerPoint deck into Metis's memory, with provenance back to the file.

    Stores the deck's content so it can be recalled and cross-referenced later —
    which is what closes the round trip: a deck edited in PowerPoint becomes
    something Metis knows about again.

    Excel files are NOT ingested. Their structure can be read (see
    read_office_document) but their contents must not enter memory — they are data,
    not authored content, and may hold patient records.

    Args:
        path: Path to a .pptx file.
        project_id: Optional project to attach it to.
    """
    p = _resolve(path)
    if not p:
        return [TextContent(type="text", text=f"No file at '{path}'.")]
    if p.suffix.lower() != ".pptx":
        return [TextContent(type="text", text=(
            "Only PowerPoint decks are ingested. A spreadsheet's contents must not enter "
            "memory — read its structure with read_office_document, and analyse the values "
            "with the safe-analysis procedure instead."
        ))]
    try:
        data = _read_pptx(p)
    except Exception as exc:
        return [TextContent(type="text", text=f"Could not read {p.name}: {exc}")]

    summary = " · ".join(s["title"] for s in data["slides"][:12])
    content = "\n\n".join(
        f"{s['number']}. {s['title']}\n{s['body']}" + (f"\nNotes: {s['notes']}" if s["notes"] else "")
        for s in data["slides"]
    )
    now = datetime.datetime.now().isoformat(timespec="seconds")
    try:
        with connect(paths.db) as con:
            con.execute(
                "CREATE TABLE IF NOT EXISTS office_documents ("
                " id INTEGER PRIMARY KEY AUTOINCREMENT,"
                " path TEXT UNIQUE, kind TEXT, title TEXT, slide_count INTEGER,"
                " summary TEXT, content TEXT, project_id TEXT,"
                " file_mtime TEXT, ingested_at TEXT)"
            )
            con.execute(
                "INSERT INTO office_documents "
                "(path, kind, title, slide_count, summary, content, project_id, file_mtime, ingested_at) "
                "VALUES (?,?,?,?,?,?,?,?,?) "
                "ON CONFLICT(path) DO UPDATE SET "
                " title=excluded.title, slide_count=excluded.slide_count, summary=excluded.summary,"
                " content=excluded.content, file_mtime=excluded.file_mtime, ingested_at=excluded.ingested_at",
                (str(p), "pptx", data["slides"][0]["title"] if data["slides"] else p.stem,
                 data["slide_count"], summary[:1000], content[:200_000], project_id,
                 datetime.datetime.fromtimestamp(p.stat().st_mtime).isoformat(timespec="seconds"),
                 now),
            )
            con.commit()
    except Exception as exc:
        return [TextContent(type="text", text=f"Read {p.name} but could not store it: {exc}")]

    # Make it findable by meaning, not just by name — the same reasoning as P3.10.
    try:
        from metis_mcp.tools.ideas import _embed_episodic

        _embed_episodic(f"PowerPoint deck '{p.name}': {summary}", "document")
    except Exception:
        pass

    return [TextContent(type="text", text=(
        f"Ingested **{p.name}** — {data['slide_count']} slide(s), stored with provenance back to "
        f"`{p}`.\nSlides: {summary[:300]}"
    ))]

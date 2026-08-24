"""authoring.py — let Metis MAKE things, not only remember them.

WHY THIS MODULE EXISTS
    An audit on 2026-08-21 asked a blunt question: what can a researcher who works
    entirely in Claude Desktop actually not do? The answer had a clean shape.
    Everything that goes through the DATABASE worked — projects, tasks, courses,
    library, memory, meetings. Everything that touches the FILESYSTEM did not.

    Of 236 tools, exactly three could write a file and all three wrote config.
    Zero could produce a document. `office.py` reads .pptx and .docx and cannot
    create either. So Metis could discuss a presentation at length and never make
    one.

    That gap was invisible in Claude Code, because Claude Code supplies Write,
    Edit and Bash itself — the assistant filled it without anyone noticing the
    tools had not. In Desktop, and in Cursor, and in whatever MCP client a future
    user picks, nothing fills it.

    These are MCP tools, so they work in every client at once. That is the point:
    a capability added here reaches Desktop, Code, and anything else that speaks
    the protocol, rather than being another Claude Code privilege.

WHERE FILES GO, AND WHY IT IS CONSTRAINED
    Everything is written under `outputs/` inside the Research Cortex, and paths
    are resolved and containment-checked before writing. A tool that can write
    anywhere is a tool that can overwrite anything, and this one is invoked by a
    model rather than a person. The constraint is not distrust of the model; it is
    that an accident here costs the researcher work they cannot recover.
"""
from __future__ import annotations

import logging
import re
import unicodedata
from datetime import datetime
from pathlib import Path

from mcp.types import TextContent

from metis_mcp.app_instance import app
from metis_mcp.config import paths

log = logging.getLogger("metis")

_MAX_SLIDES = 60
_MAX_CHARS = 400_000


def _slug(text: str, limit: int = 60) -> str:
    t = unicodedata.normalize("NFKD", text or "").encode("ascii", "ignore").decode()
    t = re.sub(r"[^A-Za-z0-9]+", "-", t).strip("-")
    return (t[:limit].strip("-") or "untitled").lower()


def _safe_target(subdir: str, filename: str) -> Path:
    """Resolve a path under outputs/ and refuse anything that escapes it."""
    base = (paths.root / "outputs" / subdir).resolve()
    base.mkdir(parents=True, exist_ok=True)
    target = (base / filename).resolve()
    if not str(target).startswith(str(base)):
        raise ValueError("refusing to write outside outputs/")
    return target


# ---------------------------------------------------------------------------
# Presentations
# ---------------------------------------------------------------------------

@app.tool()
async def create_presentation(
    title: str,
    slides: str,
    subtitle: str = "",
    filename: str = "",
) -> list[TextContent]:
    """Create a PowerPoint (.pptx) file from a simple text outline.

    The tool that closes the largest gap in Claude Desktop: Metis could read
    presentations but never produce one, so any "make me a deck" request had no
    way to land as a file.

    Args:
        title: Deck title, used on the title slide and in the filename.
        slides: The deck as plain text. Each slide starts with a line beginning
            "# " (the slide title); every following line becomes a bullet.
            A line beginning "> " becomes speaker notes for that slide.
            Blank lines are ignored. Example:

                # Why passive surveillance underperforms
                Health-facility detection depends on care-seeking
                Median delay to diagnosis exceeds 12 months
                > Open with the Mitashi 2014 figure

                # What changes with RDTs at the periphery
                Sensitivity holds at lower prevalence
        subtitle: Optional subtitle for the title slide.
        filename: Optional filename; derived from the title when omitted.

    Returns:
        The path written and a slide count.
    """
    if not title.strip():
        return [TextContent(type="text", text="A title is required.")]
    if len(slides) > _MAX_CHARS:
        return [TextContent(type="text", text="Outline too large.")]

    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return [TextContent(type="text", text=(
            "python-pptx is not installed in the Metis environment.\n"
            "Install it with:  ~/.local/share/metis-mcp/.venv/bin/pip install python-pptx"
        ))]

    # ── Parse the outline ───────────────────────────────────────────────────
    parsed: list[dict] = []
    for raw in slides.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.lstrip().startswith("# "):
            parsed.append({"title": line.lstrip()[2:].strip(),
                           "bullets": [], "notes": []})
        elif line.lstrip().startswith("> "):
            if parsed:
                parsed[-1]["notes"].append(line.lstrip()[2:].strip())
        else:
            if not parsed:
                # Bullets before any heading would otherwise be silently dropped.
                parsed.append({"title": title, "bullets": [], "notes": []})
            parsed[-1]["bullets"].append(line.strip().lstrip("-•* ").strip())

    if not parsed:
        return [TextContent(type="text", text=(
            "No slides found. Each slide must start with a line beginning '# '."
        ))]
    if len(parsed) > _MAX_SLIDES:
        return [TextContent(type="text",
                            text=f"{len(parsed)} slides exceeds the {_MAX_SLIDES} cap.")]

    # ── Build ───────────────────────────────────────────────────────────────
    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title.strip()
    if subtitle and len(s.placeholders) > 1:
        s.placeholders[1].text = subtitle.strip()

    body_layout = prs.slide_layouts[1]
    for item in parsed:
        sl = prs.slides.add_slide(body_layout)
        sl.shapes.title.text = item["title"][:120]
        body = sl.placeholders[1].text_frame
        body.clear()
        bullets = item["bullets"] or [""]
        for i, b in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = b[:300]
            para.level = 0
            for run in para.runs:
                run.font.size = Pt(20)
        if item["notes"]:
            sl.notes_slide.notes_text_frame.text = "\n".join(item["notes"])[:2000]

    stem = _slug(filename or title)
    target = _safe_target("presentations",
                          f"{datetime.now().date().isoformat()}_{stem}.pptx")
    prs.save(str(target))

    rel = target.relative_to(paths.root)
    return [TextContent(type="text", text=(
        f"Created {len(parsed) + 1} slides (title + {len(parsed)}).\n"
        f"Saved to: {rel}\n\n"
        f"Open it from the Research Cortex, or say the word and it can be revised "
        f"and rewritten to the same path."
    ))]


# ---------------------------------------------------------------------------
# Documents
# ---------------------------------------------------------------------------

@app.tool()
async def create_document(
    title: str,
    content: str,
    subdir: str = "reviews",
    fmt: str = "md",
) -> list[TextContent]:
    """Write a document to outputs/ — Markdown or Word (.docx).

    The general "save this properly" tool. Without it, anything Metis produces in
    a chat client exists only in the transcript, which is not somewhere work can
    be found again.

    Args:
        title: Document title; heads the file and derives its name.
        content: Body text. Markdown headings and bullets are honoured in both
            formats.
        subdir: Folder under outputs/ (default "reviews").
        fmt: "md" (default) or "docx".

    Returns:
        The path written.
    """
    if not title.strip() or not content.strip():
        return [TextContent(type="text", text="Both a title and content are required.")]
    if len(content) > _MAX_CHARS:
        return [TextContent(type="text", text="Content too large.")]
    if fmt not in ("md", "docx"):
        return [TextContent(type="text", text="fmt must be 'md' or 'docx'.")]

    stem = f"{datetime.now().date().isoformat()}_{_slug(title)}"

    if fmt == "md":
        target = _safe_target(subdir, f"{stem}.md")
        target.write_text(f"# {title.strip()}\n\n{content}\n", encoding="utf-8")
    else:
        try:
            from docx import Document
        except ImportError:
            return [TextContent(type="text", text=(
                "python-docx is not installed.\n"
                "Install with: ~/.local/share/metis-mcp/.venv/bin/pip install python-docx"
            ))]
        doc = Document()
        doc.add_heading(title.strip(), level=0)
        for raw in content.splitlines():
            line = raw.rstrip()
            if not line.strip():
                continue
            stripped = line.lstrip()
            if stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith(("- ", "* ", "• ")):
                doc.add_paragraph(stripped[2:], style="List Bullet")
            else:
                doc.add_paragraph(line)
        target = _safe_target(subdir, f"{stem}.docx")
        doc.save(str(target))

    return [TextContent(type="text", text=(
        f"Saved to: {target.relative_to(paths.root)}"
    ))]


@app.tool()
async def list_outputs(subdir: str = "", limit: int = 30) -> list[TextContent]:
    """List what Metis has produced, newest first.

    Exists because a file written by a tool is invisible to a chat client — the
    researcher has no way to see what was created without going to the folder.
    """
    base = (paths.root / "outputs").resolve()
    root = (base / subdir).resolve() if subdir else base
    if not str(root).startswith(str(base)) or not root.is_dir():
        return [TextContent(type="text", text=f"No such folder under outputs/: {subdir}")]

    files = sorted((p for p in root.rglob("*") if p.is_file()
                    and not p.name.startswith(".")),
                   key=lambda p: p.stat().st_mtime, reverse=True)[:max(1, limit)]
    if not files:
        return [TextContent(type="text", text=f"Nothing in outputs/{subdir}.")]

    lines = [f"{len(files)} file(s) in outputs/{subdir or ''}:", ""]
    for f in files:
        kb = f.stat().st_size // 1024
        when = datetime.fromtimestamp(f.stat().st_mtime).strftime("%Y-%m-%d %H:%M")
        lines.append(f"  {when}  {kb:>6} KB  {f.relative_to(base)}")
    return [TextContent(type="text", text="\n".join(lines))]

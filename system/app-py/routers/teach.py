"""
routers/teach.py — Teach tab: course production surface.

Builds slide outlines and PowerPoint presentations from the knowledge base.
Uses python-pptx for PPTX generation; falls back to HTML outline if unavailable.
"""

import asyncio
import datetime
import json
import os
import sqlite3
import tempfile
from pathlib import Path

import httpx
from fastapi import APIRouter, Form, Request
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.templating import Jinja2Templates

from db import db_execute, db_query, db_scalar
from models import model_for

# Resolved from system/config/models.yaml — primary + fallback chain.
# Update the YAML to swap models without code changes.
_CLAUDE_MODEL = model_for("slides")


def _get_api_key() -> str:
    return os.environ.get("ANTHROPIC_API_KEY", "")

router = APIRouter()
templates = Jinja2Templates(
    directory=str(Path(__file__).parent.parent / "templates")
)

_BUILDS_DIR = Path(os.environ.get("METIS_RC_ROOT", "")) / "outputs" / "teach-builds"
_BUILDS_DIR.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# Full-page route
# ---------------------------------------------------------------------------

@router.get("/tab/teach", response_class=HTMLResponse)
async def teach_tab(request: Request):
    return templates.TemplateResponse(request, "teach.html", {"active_tab": "teach"})


@router.get("/api/tab/teach", response_class=HTMLResponse)
async def teach_tab_partial(request: Request):
    return templates.TemplateResponse(request, "teach.html", {"active_tab": "teach"})


# ---------------------------------------------------------------------------
# Partials
# ---------------------------------------------------------------------------

@router.get("/api/teach/status")
async def teach_status():
    return JSONResponse({"ai_available": bool(_get_api_key())})


@router.get("/api/partial/teach/past-builds", response_class=HTMLResponse)
async def teach_past_builds(request: Request):
    builds = _list_builds(limit=8)
    if not builds:
        return HTMLResponse(
            '<div id="teach-past-builds" class="panel panel-pad">'
            '<div style="font-family:var(--m-mono);font-size:10px;color:var(--m-muted);">'
            'No builds yet. Use the builder to create your first course.</div></div>'
        )
    rows_html = ""
    for b in builds:
        fmt_chip = f'<span style="font-family:var(--m-mono);font-size:8px;letter-spacing:0.1em;padding:2px 6px;border:1px solid var(--m-line);border-radius:2px;color:var(--m-muted);">{b["format"].upper()}</span>'
        dl = ""
        if b.get("pptx_path") and Path(b["pptx_path"]).exists():
            dl = f'<a href="/api/teach/download?path={b["pptx_path"]}" style="font-family:var(--m-mono);font-size:9px;color:var(--m-accent);" download>↓ pptx</a>'
        rows_html += f"""
        <div style="padding:12px 16px;border-bottom:1px solid var(--m-rule-soft);">
          <div style="display:flex;align-items:baseline;gap:8px;margin-bottom:3px;">
            <span style="font-family:var(--m-display);font-size:13px;color:var(--m-ink);">{b['title'][:55]}</span>
            {fmt_chip}
          </div>
          <div style="display:flex;gap:10px;align-items:center;">
            <span style="font-family:var(--m-mono);font-size:9px;color:var(--m-muted);">{b['audience']} · {b['depth']} · {b['date']}</span>
            {dl}
          </div>
        </div>"""
    return HTMLResponse(f'<div id="teach-past-builds" class="panel" style="overflow:hidden;">{rows_html}</div>')


@router.get("/api/partial/teach/kb-stats", response_class=HTMLResponse)
async def teach_kb_stats(request: Request):
    papers  = db_scalar("SELECT COUNT(*) FROM literature_metadata", default=0) or 0
    cards   = db_scalar("SELECT COUNT(*) FROM library_cards", default=0) or 0
    ft      = db_scalar("SELECT COUNT(*) FROM library_fulltext WHERE word_count > 0", default=0) or 0
    mtgs    = db_scalar("SELECT COUNT(*) FROM meetings", default=0) or 0
    return HTMLResponse(
        f'<div id="teach-kb-stats-panel" class="panel panel-pad">'
        f'<div style="display:grid;grid-template-columns:1fr 1fr;gap:8px 16px;">'
        f'<div><div style="font-family:var(--m-mono);font-size:9px;color:var(--m-muted);">PAPERS</div>'
        f'<div style="font-family:var(--m-display);font-size:22px;color:var(--m-ink);">{papers}</div></div>'
        f'<div><div style="font-family:var(--m-mono);font-size:9px;color:var(--m-muted);">FULL TEXT</div>'
        f'<div style="font-family:var(--m-display);font-size:22px;color:var(--m-ink);">{ft}</div></div>'
        f'<div><div style="font-family:var(--m-mono);font-size:9px;color:var(--m-muted);">NOTES/IDEAS</div>'
        f'<div style="font-family:var(--m-display);font-size:22px;color:var(--m-ink);">{cards}</div></div>'
        f'<div><div style="font-family:var(--m-mono);font-size:9px;color:var(--m-muted);">MEETINGS</div>'
        f'<div style="font-family:var(--m-display);font-size:22px;color:var(--m-ink);">{mtgs}</div></div>'
        f'</div></div>'
    )


# ---------------------------------------------------------------------------
# Build endpoint
# ---------------------------------------------------------------------------

@router.post("/api/teach/build")
async def teach_build(
    request: Request,
    title:    str = Form(...),
    audience: str = Form("practitioner"),
    depth:    str = Form("overview"),
    context:  str = Form(""),
    sources:  str = Form('["library","literature","meetings"]'),
    format:   str = Form("pptx"),
):
    src_list = json.loads(sources) if sources else ["library", "literature"]
    n_slides = {"overview": 10, "lecture": 22, "workshop": 35}.get(depth, 10)

    # Pull knowledge context from DB
    kb_context = _gather_context(title, src_list, max_items=12)

    # Generate slide outline — try AI first, fall back to templates
    slides = await _generate_outline_ai(title, audience, depth, n_slides, context, kb_context)
    ai_used = slides is not None
    if not ai_used:
        slides = _generate_outline(title, audience, depth, n_slides, context, kb_context)

    # Save build record
    today = datetime.date.today().isoformat()
    slug = title.lower().replace(" ", "-")[:30]
    pptx_path = None

    if format == "pptx":
        pptx_path = _build_pptx(slides, title, audience, today, slug)

    # Save to builds index
    _save_build({
        "title": title, "audience": audience, "depth": depth,
        "date": today, "format": format,
        "pptx_path": str(pptx_path) if pptx_path else "",
        "sources": src_list,
    })

    # Render HTML preview
    html = _render_outline_html(slides)

    return JSONResponse({
        "html": html,
        "pptx_path": str(pptx_path) if pptx_path else "",
        "slide_count": len(slides),
        "ai_generated": ai_used,
    })


# ---------------------------------------------------------------------------
# Download
# ---------------------------------------------------------------------------

@router.get("/api/teach/download")
async def teach_download(path: str):
    p = Path(path)
    if not p.exists() or not str(p).startswith(str(_BUILDS_DIR)):
        return HTMLResponse("File not found", status_code=404)
    return FileResponse(str(p), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        filename=p.name)


# ---------------------------------------------------------------------------
# Knowledge context gathering
# ---------------------------------------------------------------------------

def _gather_context(title: str, sources: list[str], max_items: int = 12) -> list[dict]:
    keywords = [w for w in title.lower().split() if len(w) > 3]
    if not keywords:
        return []

    items = []
    like_clause = " OR ".join(["title LIKE ?" for _ in keywords])
    params = [f"%{k}%" for k in keywords]

    if "literature" in sources:
        rows = db_query(
            f"SELECT title, authors, year, abstract FROM literature_metadata "
            f"WHERE ({like_clause}) AND (abstract IS NOT NULL AND abstract != '') "
            f"ORDER BY year DESC LIMIT {max_items}",
            params,
        ) or []
        for r in rows:
            items.append({"type": "paper", "title": r.get("title",""), "detail": r.get("abstract","")[:200]})

    if "library" in sources:
        rows = db_query(
            f"SELECT title, summary FROM library_cards WHERE ({like_clause}) LIMIT 6",
            params,
        ) or []
        for r in rows:
            items.append({"type": "note", "title": r.get("title",""), "detail": r.get("summary","")[:150]})

    if "meetings" in sources:
        rows = db_query(
            f"SELECT title, decisions FROM meetings WHERE ({like_clause}) LIMIT 4",
            params,
        ) or []
        for r in rows:
            items.append({"type": "meeting", "title": r.get("title",""), "detail": r.get("decisions","")[:150]})

    return items[:max_items]


# ---------------------------------------------------------------------------
# Slide outline generation (template-based; Claude not called here)
# ---------------------------------------------------------------------------

_DEPTH_STRUCTURES = {
    "overview": [
        "Title & learning objectives",
        "Background and context",
        "Epidemiology / burden of disease",
        "Key concepts and definitions",
        "Current evidence — what we know",
        "Methods and approaches",
        "Key findings from the literature",
        "Challenges and gaps",
        "Conclusions and take-home messages",
        "References",
    ],
    "lecture": [
        "Title, objectives & outline",
        "Introduction — why this matters",
        "Historical background",
        "Epidemiology — global burden",
        "Epidemiology — regional patterns",
        "Aetiology and pathophysiology",
        "Clinical presentation",
        "Diagnostic approaches",
        "Treatment and management",
        "Prevention and control",
        "Surveillance and monitoring",
        "Key evidence — landmark studies",
        "Current debates and controversies",
        "Equity and access dimensions",
        "Economic and health systems considerations",
        "Case study / example",
        "Future directions",
        "Implications for practice",
        "Conclusions",
        "Discussion questions",
        "Key references",
        "Further reading",
    ],
    "workshop": [
        "Title & workshop overview",
        "Learning objectives",
        "Pre-assessment question",
        "Part 1: Background & context",
        "Epidemiological overview",
        "Burden of disease — data and trends",
        "Risk factors and determinants",
        "Part 2: Evidence base",
        "Diagnostic approaches — principles",
        "Diagnostic approaches — evidence",
        "Treatment options — overview",
        "Treatment evidence — RCTs and observational",
        "Control programmes — global",
        "Control programmes — regional",
        "Part 3: Methods",
        "Study design considerations",
        "Surveillance methods",
        "Data quality and bias",
        "Analytical approaches",
        "Part 4: Applied",
        "Case study 1",
        "Case study 2 — discussion",
        "Policy implications",
        "Implementation challenges",
        "Health economics perspective",
        "Equity dimensions",
        "Part 5: Synthesis",
        "Knowledge gaps",
        "Research priorities",
        "Summary and conclusions",
        "Key take-home messages",
        "Group exercise",
        "Evaluation and feedback",
        "References",
        "Further resources",
    ],
}

_AUDIENCE_NOTES = {
    "expert": "Assume deep domain expertise. Use technical terminology. Focus on nuance, evidence quality, and research gaps.",
    "practitioner": "Practical focus. Translate evidence into clinical/operational guidance. Avoid excessive jargon.",
    "student": "Build from first principles. Define all terms. Include learning objectives and self-assessment.",
    "public": "Plain language. Use analogies. Focus on key messages and personal relevance.",
}


async def _generate_outline_ai(title: str, audience: str, depth: str, n_slides: int,
                               context: str, kb_items: list[dict]) -> list[dict] | None:
    """Call Claude to generate slide content. Returns None on any failure."""
    api_key = _get_api_key()
    if not api_key:
        return None

    structure = _DEPTH_STRUCTURES.get(depth, _DEPTH_STRUCTURES["overview"])[:n_slides]
    aud_note = _AUDIENCE_NOTES.get(audience, "")

    kb_block = ""
    if kb_items:
        kb_block = "\n\nKnowledge base excerpts:\n" + "\n".join(
            f"- [{item['type']}] {item['title']}: {item['detail'][:120]}"
            for item in kb_items[:8]
        )

    context_block = f"\nUser-provided context: {context}" if context else ""

    system = (
        "You are an expert academic course designer. "
        "Generate concise, evidence-based slide content for a course presentation."
    )
    user_msg = (
        f"Topic: {title}\nAudience: {audience} — {aud_note}\nDepth: {depth}{context_block}{kb_block}\n\n"
        f"Generate content for these {len(structure)} slides. "
        f"For each slide respond with EXACTLY this format (one block per slide):\n"
        f"SLIDE N\n"
        f"HEADING: [exact heading text]\n"
        f"BULLETS:\n- bullet 1\n- bullet 2\n- bullet 3\n"
        f"NOTE: [one sentence speaker note]\n---\n\n"
        f"Slides:\n" + "\n".join(f"{i+1}. {h}" for i, h in enumerate(structure))
    )

    try:
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": api_key,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json",
                },
                json={
                    "model": _CLAUDE_MODEL,
                    "max_tokens": 4096,
                    "system": system,
                    "messages": [{"role": "user", "content": user_msg}],
                },
            )
        if resp.status_code != 200:
            return None
        data = resp.json()
        raw = data["content"][0]["text"]
    except Exception:
        return None

    # Parse the structured response
    slides = []
    blocks = [b.strip() for b in raw.split("---") if b.strip()]
    for i, (heading, block) in enumerate(zip(structure, blocks), 1):
        bullets = []
        note = ""
        lines = block.splitlines()
        in_bullets = False
        for line in lines:
            line = line.strip()
            if line.startswith("BULLETS:"):
                in_bullets = True
            elif line.startswith("NOTE:"):
                note = line[5:].strip()
                in_bullets = False
            elif in_bullets and line.startswith("- "):
                bullets.append(line[2:])
        if not bullets:
            bullets = _default_bullets(heading, title, audience)
        slides.append({"n": i, "heading": heading, "bullets": bullets[:5], "notes": note})

    # Pad with template slides if parse returned fewer than expected
    while len(slides) < len(structure):
        i = len(slides) + 1
        h = structure[i - 1]
        slides.append({"n": i, "heading": h, "bullets": _default_bullets(h, title, audience), "notes": ""})

    return slides


def _generate_outline(title: str, audience: str, depth: str, n_slides: int,
                      context: str, kb_items: list[dict]) -> list[dict]:
    structure = _DEPTH_STRUCTURES.get(depth, _DEPTH_STRUCTURES["overview"])[:n_slides]
    aud_note = _AUDIENCE_NOTES.get(audience, "")
    slides = []
    for i, heading in enumerate(structure, 1):
        # Find relevant KB item for this slide
        notes = ""
        for item in kb_items:
            if any(w in heading.lower() for w in item["title"].lower().split() if len(w) > 3):
                notes = item["detail"][:120]
                break
        slides.append({
            "n": i,
            "heading": heading,
            "bullets": _default_bullets(heading, title, audience),
            "notes": notes,
        })

    if context:
        slides[0]["bullets"].append(f"Context note: {context[:100]}")

    return slides


def _default_bullets(heading: str, topic: str, audience: str) -> list[str]:
    h = heading.lower()
    if "title" in h or "objective" in h:
        return [
            f"Course: {topic}",
            "After this session you will be able to:",
            "— Describe the key concepts",
            "— Interpret the evidence",
            "— Apply findings to your context",
        ]
    if "background" in h or "introduction" in h or "context" in h:
        return ["Why this topic matters", "Scope and definitions", "Key stakeholders and systems"]
    if "epidemiology" in h or "burden" in h:
        return ["Global burden estimates", "Regional distribution", "Trend over time", "High-risk groups"]
    if "method" in h or "diagnostic" in h:
        return ["Principle", "Evidence base", "Sensitivity / specificity", "Implementation considerations"]
    if "treatment" in h or "management" in h:
        return ["First-line options", "Evidence from trials", "Practical considerations", "Monitoring outcomes"]
    if "conclusion" in h or "take-home" in h or "summary" in h:
        return ["Key message 1", "Key message 2", "Key message 3", "Questions for reflection"]
    if "reference" in h:
        return [f"Key references for {topic}", "Additional reading", "Data sources"]
    return ["[Add content]", "[Add content]", "[Add content]"]


# ---------------------------------------------------------------------------
# HTML outline renderer
# ---------------------------------------------------------------------------

def _render_outline_html(slides: list[dict]) -> str:
    parts = []
    for s in slides:
        bullets_html = "".join(
            f'<div style="font-family:var(--m-sans);font-size:13px;color:var(--m-ink);padding:2px 0 2px 12px;border-left:2px solid var(--m-rule-soft);">{b}</div>'
            for b in s["bullets"]
        )
        notes_html = ""
        if s.get("notes"):
            notes_html = f'<div style="font-family:var(--m-mono);font-size:9px;color:var(--m-muted);margin-top:6px;padding:6px 8px;background:var(--m-surface-2);border-radius:3px;">KB: {s["notes"]}</div>'
        parts.append(
            f'<div style="padding:12px 0;border-bottom:1px solid var(--m-rule-soft);">'
            f'<div style="display:flex;align-items:baseline;gap:10px;margin-bottom:6px;">'
            f'<span style="font-family:var(--m-mono);font-size:9px;color:var(--m-muted);width:20px;flex-shrink:0;">{s["n"]:02d}</span>'
            f'<span style="font-family:var(--m-display);font-size:14px;font-weight:500;color:var(--m-ink);">{s["heading"]}</span>'
            f'</div>'
            f'{bullets_html}{notes_html}</div>'
        )
    count = len(slides)
    header = (
        f'<div style="margin-bottom:16px;padding-bottom:12px;border-bottom:1px solid var(--m-rule-soft);">'
        f'<span style="font-family:var(--m-mono);font-size:10px;color:var(--m-muted);">{count} SLIDES · OUTLINE PREVIEW</span>'
        f'</div>'
    )
    return header + "".join(parts)


# ---------------------------------------------------------------------------
# PPTX builder
# ---------------------------------------------------------------------------

def _build_pptx(slides: list[dict], title: str, audience: str, today: str, slug: str) -> Path | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return None

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colors
    BG     = RGBColor(0xFA, 0xF9, 0xF6)
    INK    = RGBColor(0x1A, 0x1A, 0x1A)
    ACCENT = RGBColor(0x5C, 0x5F, 0xE4)
    MUTED  = RGBColor(0x8E, 0x92, 0x9E)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

    blank_layout = prs.slide_layouts[6]  # completely blank

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)

        # Background
        bg = slide.background.fill
        if i == 0:
            bg.solid()
            bg.fore_color.rgb = ACCENT
        else:
            bg.solid()
            bg.fore_color.rgb = BG

        is_title_slide = (i == 0)
        txt_color = WHITE if is_title_slide else INK

        # Slide number (non-title)
        if not is_title_slide:
            num_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(0.5), Inches(0.4))
            tf = num_box.text_frame
            tf.text = f"{slide_data['n']:02d}"
            tf.paragraphs[0].runs[0].font.size = Pt(8)
            tf.paragraphs[0].runs[0].font.color.rgb = MUTED
            tf.paragraphs[0].runs[0].font.name = "Courier New"

        # Heading
        head_top  = Inches(0.7) if not is_title_slide else Inches(2.5)
        head_left = Inches(0.9)
        head_w    = Inches(11.5)
        head_h    = Inches(1.0)
        head_box  = slide.shapes.add_textbox(head_left, head_top, head_w, head_h)
        tf = head_box.text_frame
        tf.word_wrap = True
        tf.text = slide_data["heading"]
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(28) if is_title_slide else Pt(22)
        run.font.bold = True
        run.font.color.rgb = WHITE if is_title_slide else ACCENT
        run.font.name = "Helvetica Neue" if not is_title_slide else "Helvetica Neue"

        # Subtitle line on title slide
        if is_title_slide:
            sub_box = slide.shapes.add_textbox(head_left, Inches(3.35), head_w, Inches(0.5))
            tf2 = sub_box.text_frame
            tf2.text = f"{audience.title()} · {today}"
            tf2.paragraphs[0].runs[0].font.size = Pt(12)
            tf2.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)
            tf2.paragraphs[0].runs[0].font.name = "Helvetica Neue"

        # Divider line (non-title)
        if not is_title_slide:
            from pptx.util import Emu
            line = slide.shapes.add_connector(1, Inches(0.9), Inches(1.85), Inches(12.4), Inches(1.85))
            line.line.color.rgb = ACCENT
            line.line.width = Pt(0.75)

        # Bullets (non-title)
        if not is_title_slide and slide_data.get("bullets"):
            bul_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.5))
            tf = bul_box.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(slide_data["bullets"]):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.space_before = Pt(4)
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(14)
                run.font.color.rgb = INK
                run.font.name = "Helvetica Neue"

        # Speaker notes
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

    out_path = _BUILDS_DIR / f"{today}_{slug}.pptx"
    prs.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# Builds index
# ---------------------------------------------------------------------------

_BUILDS_INDEX = _BUILDS_DIR / "builds_index.json"


def _list_builds(limit: int = 8) -> list[dict]:
    if not _BUILDS_INDEX.exists():
        return []
    try:
        data = json.loads(_BUILDS_INDEX.read_text())
        return data[-limit:][::-1]
    except Exception:
        return []


def _save_build(build: dict):
    builds = []
    if _BUILDS_INDEX.exists():
        try:
            builds = json.loads(_BUILDS_INDEX.read_text())
        except Exception:
            builds = []
    builds.append(build)
    _BUILDS_INDEX.write_text(json.dumps(builds, indent=2))


# ---------------------------------------------------------------------------
# Gap analysis — what topics aren't covered by the library yet
# ---------------------------------------------------------------------------


@router.get("/api/partial/teach/gap-analysis", response_class=HTMLResponse)
async def teach_gap_analysis(request: Request):
    """For each active course, compare its topic keywords against the library."""
    courses = db_query(
        "SELECT id, title, slug, category FROM learning_courses "
        "WHERE status = 'active' ORDER BY title LIMIT 3",
        default=[],
    ) or []

    # Check if course_topics table is present and seeded
    topics_table = db_scalar(
        "SELECT name FROM sqlite_master WHERE type='table' AND name='course_topics'",
        default=None,
    )

    course_rows = []
    for c in courses:
        keywords: list[str] = []
        if topics_table:
            kws = db_query(
                "SELECT keyword FROM course_topics WHERE course_id = ? LIMIT 30",
                (c.get("id"),),
                default=[],
            ) or []
            keywords = [k["keyword"] for k in kws if k.get("keyword")]
        # Fall back to title words when no topics are seeded
        if not keywords:
            keywords = [w for w in (c.get("title") or "").split() if len(w) > 3][:8]

        covered: list[str] = []
        gaps: list[str] = []
        for kw in keywords:
            like = f"%{kw}%"
            n = db_scalar(
                "SELECT COUNT(*) FROM literature_metadata WHERE title LIKE ? OR abstract LIKE ?",
                (like, like),
                default=0,
            ) or 0
            if n > 0:
                covered.append(f"{kw} ({n})")
            else:
                gaps.append(kw)

        total = len(covered) + len(gaps)
        covered_pct = int(round((len(covered) / total) * 100)) if total else 0
        gap_pct = 100 - covered_pct if total else 0

        course_rows.append({
            "title": c.get("title") or "Untitled course",
            "category": c.get("category") or "",
            "covered": covered,
            "gaps": gaps,
            "top_gaps": gaps[:3],
            "covered_n": len(covered),
            "gaps_n": len(gaps),
            "total": total,
            "covered_pct": covered_pct,
            "gap_pct": gap_pct,
        })

    return templates.TemplateResponse(
        request,
        "partials/teach_gap_analysis.html",
        {"courses": course_rows},
    )


# ---------------------------------------------------------------------------
# Q-bank — placeholder until question generation lands
# ---------------------------------------------------------------------------


@router.get("/api/partial/teach/qbank", response_class=HTMLResponse)
async def teach_qbank(request: Request):
    """Show question counts per course, or a friendly placeholder if no Q-bank yet."""
    has_table = db_scalar(
        "SELECT name FROM sqlite_master WHERE type='table' AND name='course_questions'",
        default=None,
    )
    rows = []
    total = 0
    if has_table:
        rows = db_query(
            "SELECT course_id, COUNT(*) as n FROM course_questions GROUP BY course_id ORDER BY n DESC LIMIT 10",
            default=[],
        ) or []
        total = db_scalar("SELECT COUNT(*) FROM course_questions", default=0) or 0
        # Resolve course titles + first question text as a teaser
        for r in rows:
            t = db_query(
                "SELECT title FROM learning_courses WHERE id=? LIMIT 1",
                (r.get("course_id"),),
                default=[],
            )
            r["title"] = (t[0]["title"] if t else "(unknown course)")
            # Try a few likely column names for the question body
            sample_text = ""
            for col in ("question_text", "question", "prompt", "text"):
                try:
                    s = db_query(
                        f"SELECT {col} AS body FROM course_questions WHERE course_id=? "
                        f"ORDER BY rowid ASC LIMIT 1",
                        (r.get("course_id"),),
                        default=[],
                    ) or []
                    if s and s[0].get("body"):
                        sample_text = s[0]["body"]
                        break
                except Exception:
                    continue
            if len(sample_text) > 80:
                sample_text = sample_text[:80].rsplit(" ", 1)[0] + "..."
            r["sample"] = sample_text

    return templates.TemplateResponse(
        request,
        "partials/teach_qbank.html",
        {"rows": rows, "total": total, "has_table": bool(has_table)},
    )


# ---------------------------------------------------------------------------
# Legacy endpoints kept for backward compat (return empty gracefully)
# ---------------------------------------------------------------------------

@router.get("/api/partial/teach/meta", response_class=HTMLResponse)
async def teach_meta(request: Request):
    total = db_scalar("SELECT COUNT(*) FROM learning_courses", default=0) or 0
    return HTMLResponse(f"0 BUILDS · {total} COURSES IN LEARNING")

@router.get("/api/partial/teach/active-label", response_class=HTMLResponse)
async def _teach_active_label(request: Request):
    return HTMLResponse('<span>Course builder</span><span class="tail">READY</span>')

@router.get("/api/partial/teach/active-draft", response_class=HTMLResponse)
async def _teach_active_draft(request: Request):
    return HTMLResponse('<div></div>')

@router.get("/api/partial/teach/courses-list", response_class=HTMLResponse)
async def _teach_courses_list(request: Request):
    return HTMLResponse('<div></div>')

@router.get("/api/partial/teach/suggested", response_class=HTMLResponse)
async def _teach_suggested(request: Request):
    return HTMLResponse('<div></div>')
